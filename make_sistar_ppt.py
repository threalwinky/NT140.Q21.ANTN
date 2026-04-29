from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

ROOT = Path('/home/team/NT140.Q21.ANTN/Final')
OUT = ROOT / 'SISTAR_presentation_vi.pptx'
IMG_ACC = ROOT / 'reproduction/output/classification_accuracy.png'
IMG_THRESH = ROOT / 'reproduction/output/threshold_usage.png'
IMG_PUSH = ROOT / 'reproduction/output/pushback_attack_bytes.png'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

TITLE = RGBColor(21, 66, 139)
TEXT = RGBColor(40, 40, 40)
ACCENT = RGBColor(192, 57, 43)
BG = RGBColor(248, 250, 252)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12.1), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = 'Arial'
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = TITLE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = 'Arial'
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT


def add_bullets(slide, items, left=0.8, top=1.3, width=11.7, height=5.6, font_size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p.level = level
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(font_size - level * 2)
        p.font.color.rgb = TEXT
        if level == 0:
            p.space_after = Pt(8)


def add_two_column(slide, left_items, right_items, title, subtitle=None):
    set_bg(slide)
    add_title(slide, title, subtitle)
    add_bullets(slide, left_items, left=0.7, top=1.35, width=5.8, height=5.6, font_size=17)
    add_bullets(slide, right_items, left=6.8, top=1.35, width=5.8, height=5.6, font_size=17)


def add_image_slide(slide, title, bullets, image_path, image_left=7.0, image_top=1.5, image_width=5.5):
    set_bg(slide)
    add_title(slide, title)
    add_bullets(slide, bullets, left=0.7, top=1.35, width=5.9, height=5.5, font_size=17)
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(image_left), Inches(image_top), width=Inches(image_width))


def add_title_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.8), Inches(1.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = 'SISTAR: phát hiện và giảm thiểu DDoS\nbằng Programmable Data Plane'
    r.font.name = 'Arial'
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = TITLE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = 'Tóm tắt paper, đánh giá thực tế, và hướng dẫn hiện thực hóa một phần ý tưởng'
    p2.font.name = 'Arial'
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.text = 'Paper: ACM CCS 2025 | File: paper.pdf'
    p3.font.name = 'Arial'
    p3.font.size = Pt(14)
    p3.font.color.rgb = ACCENT


def add_section(title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, title, subtitle)
    add_bullets(slide, bullets)


add_title_slide()

add_section('1. Bài toán paper giải quyết', [
    'Phát hiện và giảm thiểu DDoS ngay trong mạng, thay vì phụ thuộc nhiều vào server hoặc control plane.',
    'Xử lý cả flooding attack và application-layer / protocol-based attack.',
    'Mục tiêu kép: giữ độ chính xác cao nhưng vẫn tiết kiệm tài nguyên trên programmable switch.',
    'Bối cảnh trọng tâm: triển khai trên nhiều switch để phòng thủ phân tán, không chỉ bảo vệ cục bộ.'
])

add_two_column(
    prs.slides.add_slide(prs.slide_layouts[6]),
    [
        'Mô hình đe dọa',
        (1, 'Kẻ tấn công bơm lưu lượng độc hại từ nhiều nguồn để làm cạn băng thông hoặc tài nguyên dịch vụ.'),
        (1, 'Lưu lượng đi qua topo kiểu gateway leaf → spine → server leaf.'),
        (1, 'Tấn công gồm cả SYN/UDP/ICMP flood và attack tầng ứng dụng như GoldenEye, Hulk, Slowloris.'),
    ],
    [
        'Giả định của paper',
        (1, 'Mạng thuộc một miền quản trị thống nhất.'),
        (1, 'Có programmable switch hỗ trợ P4/PISA.'),
        (1, 'Switch có thể trích một số packet/flow feature và gửi alert packet cho nhau.'),
        (1, 'Có thể retrain mô hình khi xuất hiện attack mới.'),
    ],
    '2. Mô hình đe dọa và giả định'
)

add_section('3. Vì sao bài toán khó?', [
    'Accuracy cao thường đòi hỏi nhiều feature và nhiều threshold, nhưng switch lại bị giới hạn SRAM, TCAM, stage và table entry.',
    'Mapping mô hình ML sang P4 pipeline không trực tiếp; càng nhiều threshold thì encoding càng dài và bảng càng lớn.',
    'Một switch đơn lẻ khó xử lý tốt DDoS phân tán, nên phải có cơ chế hợp tác giữa nhiều switch.',
    'Paper nhắm đúng trade-off giữa detection quality và deployability.'
])

add_section('4. Ý tưởng chính của SISTAR', [
    'Chọn một tập feature nhỏ nhưng giàu thông tin, kết hợp packet-level và flow-level feature.',
    'Dùng DT-CTS để giảm số threshold khi huấn luyện decision tree.',
    'Mã hóa các feature thành vector bit rồi dùng ternary match để ánh xạ sang match-action table.',
    'Triển khai mô hình phân tán: gateway dùng mô hình nhẹ, spine dùng mô hình mạnh hơn.',
    'Khi phát hiện tấn công, switch sinh alert packet và kích hoạt pushback lên upstream switch.'
])

add_section('5. DT-CTS là gì?', [
    'DT-CTS = Decision Tree with Constrained Threshold Segmentation.',
    'Khác với decision tree thường: mỗi feature bị giới hạn số threshold tối đa trong quá trình xây cây.',
    'Lợi ích: giảm node dư thừa, giảm số điều kiện so khớp, giảm độ phức tạp khi deploy lên switch.',
    'Đây là đóng góp quan trọng vì nó tối ưu mô hình theo ràng buộc phần cứng mạng, không chỉ theo accuracy.'
])

add_section('6. Cách SISTAR triển khai trên data plane', [
    'Mỗi feature được chia thành các khoảng theo threshold.',
    'Mỗi khoảng được gán một mã nhị phân ngắn.',
    'Ghép các mã thành vector đặc trưng `bin_feature`.',
    'Dùng ternary match để ánh xạ vector này sang class / action.',
    'Thiết kế này giúp giảm stage, giảm table entry và phù hợp với match-action pipeline của P4 switch.'
])

add_section('7. Cơ chế phân tán và pushback', [
    'Gateway switch phát hiện nhanh các dấu hiệu tấn công rõ ràng bằng mô hình nhẹ.',
    'Spine switch xử lý các mẫu phức tạp hơn bằng mô hình chi tiết hơn.',
    'Khi một switch nghi ngờ có tấn công, nó phát alert packet tới các switch khác.',
    'Nếu mức cảnh báo vượt ngưỡng, hệ thống pushback lên upstream switch để chặn gần nguồn hơn.',
    'Ý nghĩa: giảm tải cho nạn nhân và giảm phần băng thông độc hại đi sâu vào mạng.'
])

add_section('8. Kết quả chính trong paper', [
    'DT-CTS giảm threshold khoảng 70% mà vẫn giữ accuracy cao.',
    'Paper báo cáo chỉ cần 3 feature để đạt khoảng 98% F1.',
    'Trên traffic DDoS thực trong topo phân tán, độ chính xác vượt 95%; với traffic mô phỏng, đạt 99.7%.',
    'Pushback giúp giảm attack traffic bandwidth utilization tới khoảng 40%.',
    'Độ trễ xử lý vẫn thấp, khoảng 820–920 ns, và hệ thống giữ line-rate trên Tofino.'
])

add_section('9. Kết quả tài nguyên trong paper', [
    'SISTAR dùng khoảng 2.9% SRAM và 2.1% TCAM.',
    'tMatch xBar khoảng 1.8%.',
    'Chỉ dùng khoảng 6 stages và 7 tables.',
    'Số table entry thấp hơn đáng kể so với nhiều framework khác ở mức accuracy tương đương.',
    'Điểm mạnh lớn của paper là: mô hình không chỉ chính xác mà còn deployable trên switch thật.'
])

add_section('10. Phần nào có thể hiện thực hóa trong repo này?', [
    'Repo có reproduction rút gọn của ý tưởng SISTAR trong thư mục `reproduction/`.',
    'Có thể demo tốt 3 phần: so sánh DT/RF/DT-CTS, đo threshold-count, và mô phỏng pushback 3-hop.',
    'Repo cũng chứa P4 code trong `SISTAR/BMv2/` và `SISTAR/tofino/`, nhưng thí nghiệm hiện tại không tái hiện đầy đủ môi trường phần cứng như paper.',
    'Vì vậy, demo phù hợp nhất là bản reduced reproduction, không phải full hardware reproduction.'
])

add_section('11. Cách chạy phần hiện thực hóa', [
    'Chạy 1 lệnh: `python3 /home/team/NT140.Q21.ANTN/Final/reproduction/src/run_reproduction.py`',
    'Script sẽ tự dùng dataset CICIDS2017 Wednesday subset nếu có.',
    'Sau khi chạy xong, kết quả nằm trong `reproduction/output/`.',
    'Các file quan trọng: classification_metrics.csv, threshold_metrics.csv, pushback_metrics.csv, summary.md và các biểu đồ PNG.'
])

add_image_slide(
    prs.slides.add_slide(prs.slide_layouts[6]),
    '12. Kết quả hiện thực rút gọn: độ chính xác',
    [
        'Kết quả chạy hiện tại trên repo dùng CICIDS2017 Wednesday subset.',
        'DT: accuracy 0.9912, F1 0.9912.',
        'RF: accuracy 0.9912, F1 0.9912.',
        'DT-CTS: accuracy 0.9640, F1 0.9646.',
        'Kết luận: trong bản tái hiện rút gọn, DT-CTS đánh đổi một phần accuracy để giảm độ phức tạp triển khai.'
    ],
    IMG_ACC,
    image_left=7.0,
    image_top=1.55,
    image_width=5.6,
)

add_image_slide(
    prs.slides.add_slide(prs.slide_layouts[6]),
    '13. Kết quả hiện thực rút gọn: threshold',
    [
        'DT dùng 22 threshold tổng cộng.',
        'DT-CTS dùng 15 threshold.',
        'Giảm khoảng 31.8% so với DT trong lần chạy này.',
        'Ý nghĩa: ít threshold hơn thì dễ mã hóa hơn, ít rule hơn và gần hơn với mục tiêu deploy lên switch.'
    ],
    IMG_THRESH,
    image_left=7.0,
    image_top=1.7,
    image_width=5.4,
)

add_image_slide(
    prs.slides.add_slide(prs.slide_layouts[6]),
    '14. Kết quả hiện thực rút gọn: pushback',
    [
        'No pushback: attack bytes tới victim ≈ 1,731,491.',
        'Gated pushback: còn ≈ 51,302.',
        'Giảm khoảng 97.0% attack bytes trong mô phỏng rút gọn.',
        'Immediate pushback chặn mạnh hơn nhưng gây false block nhiều hơn.',
        'Gated pushback là cải tiến nhỏ hợp lý để cân bằng bảo vệ và nhầm lẫn.'
    ],
    IMG_PUSH,
    image_left=6.7,
    image_top=1.55,
    image_width=5.8,
)

add_two_column(
    prs.slides.add_slide(prs.slide_layouts[6]),
    [
        'Ưu điểm',
        (1, 'Tập trung đúng bài toán hệ thống: accuracy + deployability.'),
        (1, 'Tận dụng PDP để phát hiện sớm, độ trễ thấp.'),
        (1, 'Có cơ chế phân tán và pushback hợp lý.'),
        (1, 'DT-CTS gắn chặt với ràng buộc phần cứng mạng.'),
    ],
    [
        'Hạn chế',
        (1, 'Chưa có threat model hình thức thật chặt.'),
        (1, 'Phụ thuộc mạnh vào feature, dataset và topology.'),
        (1, 'Gặp attack mới vẫn cần retrain.'),
        (1, 'Có nguy cơ alert storm, forged alert, threshold inference.'),
    ],
    '15. Ưu điểm và hạn chế'
)

add_section('16. Mức độ phù hợp trong thực tế', [
    'Phù hợp cao nếu tổ chức có programmable switch thật hoặc lab P4/BMv2 và có quyền kiểm soát nhiều hop trong mạng.',
    'Phù hợp trung bình nếu hạ tầng hiện tại chưa có năng lực P4/Tofino nhưng vẫn muốn nghiên cứu giải pháp in-network defense.',
    'Ít phù hợp nếu môi trường thay đổi quá nhanh nhưng không có pipeline retraining tốt, hoặc không kiểm soát được upstream hạ tầng mạng.',
    'Trong bối cảnh thực tế, ý tưởng này hợp với data center, enterprise network, hoặc môi trường nghiên cứu có hạ tầng mạng hiện đại.'
])

add_section('17. Kết luận và thông điệp trình bày', [
    'SISTAR mạnh không chỉ vì phát hiện DDoS tốt, mà vì nó được thiết kế để chạy được trên switch thật.',
    'Đóng góp cốt lõi: DT-CTS, feature encoding, và pushback phân tán.',
    'Repo hiện tại tái hiện tốt phần lõi ý tưởng, nhưng chưa tái hiện đầy đủ thí nghiệm phần cứng của paper.',
    'Thông điệp nên chốt khi thuyết trình: paper giải quyết bài toán cân bằng giữa độ chính xác, chi phí tài nguyên, và khả năng phòng thủ phân tán.'
])

prs.save(str(OUT))
print(OUT)
